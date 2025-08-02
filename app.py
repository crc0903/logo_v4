import streamlit as st
from PIL import Image, ImageChops
import os
import io
import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

PRELOADED_LOGO_DIR = "preloaded_logos"

def load_preloaded_logos():
    logos = []
    if not os.path.exists(PRELOADED_LOGO_DIR):
        os.makedirs(PRELOADED_LOGO_DIR)
    for file in os.listdir(PRELOADED_LOGO_DIR):
        if file.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
            name = os.path.splitext(file)[0]
            image = Image.open(os.path.join(PRELOADED_LOGO_DIR, file)).convert("RGBA")
            logos.append((name.lower(), image))
    return logos

# Trims whitespace or transparency
def trim_whitespace(image):
    bg = Image.new("RGBA", image.size, (255, 255, 255, 0))
    diff = ImageChops.difference(image, bg)
    bbox = diff.getbbox()
    return image.crop(bbox) if bbox else image

# Resize image proportionally to fill a 5:2 box
def resize_to_fit_box(image, box_width_px, box_height_px, buffer_ratio=0.9):
    img_w, img_h = image.size
    img_ratio = img_w / img_h

    max_width = int(box_width_px * buffer_ratio)
    max_height = int(box_height_px * buffer_ratio)

    if img_ratio >= (max_width / max_height):
        new_width = max_width
        new_height = int(new_width / img_ratio)
    else:
        new_height = max_height
        new_width = int(new_height * img_ratio)

    return image.resize((new_width, new_height), Image.LANCZOS)

# Create a logo grid slide
def create_logo_slide(prs, logos, canvas_width_in, canvas_height_in, logos_per_row):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas_width_px = int(canvas_width_in * 96)
    canvas_height_px = int(canvas_height_in * 96)

    logo_count = len(logos)
    cols = logos_per_row if logos_per_row else max(1, round(math.sqrt(logo_count * canvas_width_in / canvas_height_in)))
    rows = math.ceil(logo_count / cols)

    cell_width = canvas_width_px / cols
    cell_height = canvas_height_px / rows

    left_margin = Inches((10 - canvas_width_in) / 2)
    top_margin = Inches((7.5 - canvas_height_in) / 2)

    for idx, (name, image) in enumerate(logos):
        col = idx % cols
        row = idx // cols

        trimmed = trim_whitespace(image)
        resized = resize_to_fit_box(trimmed, int(cell_width), int(cell_height))

        img_stream = io.BytesIO()
        resized.save(img_stream, format="PNG")
        img_stream.seek(0)

        x_offset = (cell_width - resized.width) / 2
        y_offset = (cell_height - resized.height) / 2
        left = left_margin + Inches((col * cell_width + x_offset) / 96)
        top = top_margin + Inches((row * cell_height + y_offset) / 96)

        slide.shapes.add_picture(
            img_stream, left, top,
            width=Inches(resized.width / 96),
            height=Inches(resized.height / 96)
        )

        # Optional: draw 5x2 guideline box
        guide_left = left_margin + Inches(col * cell_width / 96)
        guide_top = top_margin + Inches(row * cell_height / 96)
        guide = slide.shapes.add_shape(
            autoshape_type_id=1,  # rectangle
            left=guide_left,
            top=guide_top,
            width=Inches(cell_width / 96),
            height=Inches(cell_height / 96)
        )
        guide.line.color.rgb = RGBColor(100, 100, 100)
        guide.fill.background()  # transparent fill

# --- Streamlit UI ---
st.title("Logo Grid PowerPoint Exporter")
st.markdown("Upload logos or use preloaded ones below:")

uploaded_files = st.file_uploader("Upload logos", type=["png", "jpg", "jpeg", "webp"], accept_multiple_files=True)
preloaded = load_preloaded_logos()
selected_preloaded = st.multiselect("Select preloaded logos", options=sorted([name for name, _ in preloaded]))

canvas_width_in = st.number_input("Grid width (inches)", min_value=1.0, max_value=20.0, value=10.0)
canvas_height_in = st.number_input("Grid height (inches)", min_value=1.0, max_value=20.0, value=7.5)
logos_per_row = st.number_input("Logos per row (optional)", min_value=0, max_value=50, value=0)

if st.button("Generate PowerPoint"):
    logos = []

    if uploaded_files:
        for f in uploaded_files:
            name = os.path.splitext(f.name)[0].lower()
            image = Image.open(f).convert("RGBA")
            logos.append((name, image))

    for name, image in preloaded:
        if name in [n.lower() for n in selected_preloaded]:
            logos.append((name, image))

    if not logos:
        st.warning("Please upload or select logos.")
    else:
        # Alphabetize all logos together
        logos.sort(key=lambda x: x[0])

        prs = Presentation()
        create_logo_slide(prs, logos, canvas_width_in, canvas_height_in,
                          logos_per_row if logos_per_row > 0 else None)

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        st.success("PowerPoint created!")
        st.download_button("Download .pptx", output, file_name="logo_grid.pptx")
